import collections
import collections.abc
from pptx import Presentation
import pandas as pd

def read_ppt(filename):
    presentation = Presentation(filename)
    tables = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = ''
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                cell_text += run.text
                        row_data.append(cell_text)
                    table_data.append(row_data)
                df = pd.DataFrame(table_data)
                tables.append(df)
    return tables

tables = read_ppt('../data/PPT Reader Sample.pptx')

# Let's print the first table as an example
if tables:
    print(tables[0])
